#!/usr/bin/env python3
"""
Convert markdown pitch deck to PowerPoint presentation.
Supports slide titles, bullet points, and optional images.
"""

import argparse
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def parse_markdown(md_content):
    """Parse markdown content into slides."""
    slides = []
    current_slide = None

    # Split by slide markers (## Slide N: or # Slide N:)
    slide_pattern = r'^#{1,2}\s+(?:Slide\s+\d+:\s+)?(.+)$'

    lines = md_content.split('\n')

    for line in lines:
        line = line.rstrip()

        # Check if it's a slide title
        match = re.match(slide_pattern, line)
        if match:
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'title': match.group(1).strip(),
                'bullets': [],
                'notes': [],
                'image': None
            }
        elif current_slide:
            # Check for bullets
            if line.startswith('- ') or line.startswith('* '):
                current_slide['bullets'].append(line[2:].strip())
            # Check for numbered bullets
            elif re.match(r'^\d+\.\s+', line):
                current_slide['bullets'].append(re.sub(r'^\d+\.\s+', '', line))
            # Check for images
            elif line.startswith('!['):
                img_match = re.match(r'!\[.*?\]\((.+?)\)', line)
                if img_match:
                    current_slide['image'] = img_match.group(1)
            # Check for speaker notes (lines starting with "Note:" or "Speaker notes:")
            elif line.lower().startswith(('note:', 'speaker notes:', 'notes:')):
                current_slide['notes'].append(line.split(':', 1)[1].strip())
            # Add non-empty lines as notes if they're not separators
            elif line and line != '---' and not line.startswith('#'):
                if current_slide['bullets']:
                    # If we have bullets, treat this as a note
                    current_slide['notes'].append(line)

    if current_slide:
        slides.append(current_slide)

    return slides


def add_slide_with_content(prs, slide_data, images_dir=None, slide_number=None):
    """Add a slide to the presentation with title and bullets."""
    # Use title and content layout (index 1)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = slide_data['title']
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    # Add bullets
    if slide_data['bullets']:
        body = slide.placeholders[1]
        text_frame = body.text_frame
        text_frame.clear()

        for i, bullet in enumerate(slide_data['bullets']):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(12)

    # Add image if specified or from images_dir
    image_path = None
    if slide_data.get('image'):
        image_path = slide_data['image']
    elif images_dir and slide_number:
        # Look for images matching the slide number
        for pattern in [f"{slide_number}.png", f"{slide_number}.jpg",
                       f"slide_{slide_number}.png", f"slide_{slide_number}.jpg",
                       f"image_{slide_number}.png", f"image_{slide_number}.jpg",
                       f"{slide_number:02d}.png", f"{slide_number:02d}.jpg"]:
            potential_path = Path(images_dir) / pattern
            if potential_path.exists():
                image_path = str(potential_path)
                break

    if image_path and Path(image_path).exists():
        # Add image to the right side of the slide
        left = Inches(6)
        top = Inches(2)
        width = Inches(3.5)
        slide.shapes.add_picture(str(image_path), left, top, width=width)

    # Add speaker notes
    if slide_data['notes']:
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = '\n'.join(slide_data['notes'])


def create_title_slide(prs, title, subtitle=""):
    """Create a title slide."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 56, 100)

    if subtitle:
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(24)


def main():
    parser = argparse.ArgumentParser(description='Convert markdown pitch deck to PowerPoint')
    parser.add_argument('input', help='Input markdown file')
    parser.add_argument('-o', '--output', required=True, help='Output PowerPoint file (.pptx)')
    parser.add_argument('--images-dir', help='Directory containing images (1.png, 2.png, etc.)')
    parser.add_argument('--title', help='Override deck title')
    parser.add_argument('--subtitle', help='Deck subtitle')

    args = parser.parse_args()

    # Read markdown file
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file '{args.input}' not found")
        return 1

    with open(input_path, 'r', encoding='utf-8') as f:
        md_content = f.read()

    # Parse slides
    slides = parse_markdown(md_content)

    if not slides:
        print("Error: No slides found in markdown file")
        return 1

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create title slide
    deck_title = args.title if args.title else slides[0]['title']
    deck_subtitle = args.subtitle if args.subtitle else ""
    create_title_slide(prs, deck_title, deck_subtitle)

    # Add content slides
    for i, slide_data in enumerate(slides, 1):
        add_slide_with_content(prs, slide_data, args.images_dir, i)

    # Save presentation
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"[SUCCESS] Created PowerPoint presentation: {output_path}")
    print(f"[SUCCESS] Total slides: {len(slides) + 1} (1 title + {len(slides)} content)")

    return 0


if __name__ == '__main__':
    exit(main())
